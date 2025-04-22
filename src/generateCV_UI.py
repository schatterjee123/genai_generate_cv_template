import gradio as gr
import os
import openai

from langchain.agents import initialize_agent, AgentType
from langchain.tools import Tool
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
import docx
import pptx
import os
import ast
from pptx.util import Pt

from dotenv import load_dotenv

load_dotenv()

openai.api_key = os.getenv("OPENAI_API_KEY")

# llm_model="gpt-3.5-turbo"
llm_model="gpt-4"
# llm_model="gpt-4o-mini"

# Function to extract raw text from the DOCX CV file
def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return text

# Function that asks an LLM to structure the CV text into required sections
def llm_extract_cv_sections(cv_text):
    prompt = PromptTemplate(
        input_variables=["cv_text"],
        template="""
         Extract structured information from the following CV text:

        {cv_text}

        Format the response as JSON with the following keys:
        - NAME: Full name of the person.
        - SKILLS: A concise list of technical skills. 
        - SUMMARY: A brief summary of experience and expertise. 
        - COMPANY_1, ROLE_1, DATES_1, RESPONSIBILITIES_1: Details of the most recent job.

        Please categories the skills like Programming Language, Database etc. Please restrict to 5 skills only.

        Please limit the summary to 100 words.

        Ensure only the five most recent jobs are included. If number of jobs are less than five then extract all
        Please display the jobs information extracted in the format of COMPANY, ROLE and DATES as above.
        Which means display second job's information as COMPANY_2, ROLE_2, DATES_2, RESPONSIBILITIES_2
        Display third job's information as COMPANY_3, ROLE_3, DATES_3,, RESPONSIBILITIES_3 and so on.
        
        RESPONSIBILITIES should cover a brief summary of the project responsibilities.
        Please limit the RESPONSIBILITIES of each ROLE to 20 words.

        Provide only the JSON response with no extra text. If there is missing information, still output a valid JSON with empty values for the missing fields.

        Validate the JSON output to make sure all invalid syntax are corrected.
        """
    )
    structured_data = llm.predict(prompt.format(cv_text=cv_text))
    
    return ast.literal_eval(structured_data)  # Convert LLM response to a Python dictionary
    

# Function to fill the PowerPoint template with extracted data
def fill_ppt_template(template_path, output_path, cv_data):
    formatted_skills = ""
    for category, skills in cv_data["SKILLS"].items():
        formatted_skills += f"{category}: "  # Add category name (e.g., Programming Language)
        formatted_skills += ", ".join(skills)  # Join the skills with commas
        formatted_skills += "\n"  # Add a newline for separation between categories

    prs = pptx.Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in cv_data.items():
                    if f"{{{{{key}}}}}" in shape.text:
                        if isinstance(value, dict):
                            # value = "\n".join(f"{k}: {v}" for k, v in value.items())  # Convert dict to string
                            value = formatted_skills 
                            print("printing value:", value)
                        if isinstance(value, list):
                            value = ", ".join(value)

                        # Extract font from template placeholder before replacing text
                        template_font = None
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if f"{{{{{key}}}}}" in run.text:
                                    template_font = run.font
                                    break
                            if template_font:
                                break
                        
                        shape.text = shape.text.replace(f"{{{{{key}}}}}", value)
                        
                        # Apply extracted font settings to new text
                        if template_font:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = template_font.name
                                    run.font.size = template_font.size
                                    run.font.bold = template_font.bold
                                    run.font.italic = template_font.italic
                                    run.font.underline = template_font.underline
                                    if template_font.color and hasattr(template_font.color, "rgb"):
                                        run.font.color.rgb = template_font.color.rgb  # Handle _SchemeColor
    prs.save(output_path)
    return output_path

def extract_text_wrapper(args):
    input_dict = ast.literal_eval(args)
    docx_path = input_dict.get('docx_path')
    if not docx_path:
        raise ValueError("Missing required argument: 'docx_path'")
    return extract_text_from_docx(docx_path)

extract_text_tool = Tool(
    name="Extract CV Text",
    func=extract_text_wrapper,
    description="Extracts raw text from a CV document. Requires 'docx_path'."
)


def llm_extraction_wrapper(args):
   input = ast.literal_eval(args)
   cv_text = input.get('cv_text') 
   return llm_extract_cv_sections(cv_text)

llm_extraction_tool = Tool(
    name="LLM Extract CV Data",
    func = llm_extraction_wrapper,
    description="Uses an LLM to extract structured sections from the CV text. Requires 'cv_text'."
)


def fill_ppt_wrapper(args):
    input = ast.literal_eval(args)
    template_path = input.get('template_path') 
    output_path = input.get('output_path') 
    cv_data = input.get('cv_data') 
    return fill_ppt_template(template_path, output_path, cv_data)

fill_ppt_tool = Tool(
    name="Fill PPT Template",
    func=fill_ppt_wrapper,
    description="Fills a PowerPoint template with extracted CV data. Expects 'template_path', 'output_path', and 'cv_data'."
)


llm = ChatOpenAI(
    model_name=llm_model
)
agent = initialize_agent(
    tools=[extract_text_tool, llm_extraction_tool, fill_ppt_tool],
    llm=llm,
    agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
    verbose=True
)

# Function to handle file selection and process
def process_file(input_dir, template_file, output_dir):

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Process the CV using the agent
    for filename in os.listdir(input_dir):
        if filename.endswith(".docx"):
            input_cv_path = os.path.join(input_dir, filename)
            output_ppt_path = os.path.join(output_dir, filename.replace(".docx", "_langchain_ui.pptx"))

            agent.invoke({
                "input": {
                "docx_path": input_cv_path,
                "template_path": template_file,
                "output_path": output_ppt_path
                }})

    return f"Generated CV saved at: {output_ppt_path}."

# Create Gradio interface with a Submit button
with gr.Blocks(css="""
    .input-textbox {
        background-color: #f0f8ff;  /* Light Blue background for input fields */
        border: 2px solid #4caf50;  /* Green border */
        color: #333;  /* Dark text color */
        font-size: 16px;
    }
    .input-textbox:focus {
        background-color: #e0f7fa;  /* Lighter blue when focused */
        border: 2px solid #00796b;  /* Darker green border when focused */
    }
    .button {
        background-color: #ffeb3b;  /* Yellow background for button */
        border: 2px solid #fbc02d;  /* Darker yellow border */
        color: #000;  /* Black text color for the button */
        font-size: 16px;
    }
    .button:hover {
        background-color: #fbc02d;  /* Darken yellow when hovering */
        border: 2px solid #ff9800;  /* Darker border on hover */
    }
    .output-textbox {
        background-color: #ffffff;  /* White background for output */
        border: 2px solid #008080;  /* Teal border */
        color: #333;  /* Dark text color */
        font-size: 16px;
    }
    .title {
        font-size: 24px;
        font-weight: bold;  /* Make the title bold */
        color: #00796b;
        text-align: center;
        margin-bottom: 20px;
    }
""") as demo:
    
    gr.Markdown("<div class='title'>Riskcare CV Generator</div>")
    # Define the input components
    input_dir = gr.Textbox(
        label="Input Directory", 
        placeholder="Enter the path to the input directory of the CVs", 
        elem_classes=["input-textbox"]
    )
    template_file = gr.File(
        label="Template File", 
        file_count="single"
        # placeholder="Enter the full path of the template file", 
        # elem_classes=["input-textbox"]
    )
    output_dir = gr.Textbox(
        label="Output Directory", 
        placeholder="Enter the path to the output directory", 
        elem_classes=["input-textbox"]
    )
    
    # Define the Submit button with custom styling
    submit_btn = gr.Button("Submit", elem_classes=["button"])
    
    # Define the output component
    output = gr.Textbox(elem_classes=["output-textbox"])

    # Set the function to be triggered on button click
    submit_btn.click(fn=process_file, 
                     inputs=[input_dir, template_file, output_dir], 
                     outputs=output)



# Launch the interface
demo.launch()
